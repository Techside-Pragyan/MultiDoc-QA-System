import os
import PyPDF2
import docx
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sqlalchemy.orm import Session
from backend.database import models
from backend.vector_store.faiss_store import faiss_store
from backend.utils.config import settings

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    def extract_text(self, file_path: str, file_type: str) -> str:
        text = ""
        try:
            if file_type == 'pdf':
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            elif file_type == 'docx':
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif file_type == 'txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()
            elif file_type == 'pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting text from {file_path}: {str(e)}")
            
        return text

    def process_document(self, db: Session, document: models.Document):
        # 1. Extract text
        text = self.extract_text(document.file_path, document.file_type)
        if not text:
            return False

        # 2. Split into chunks
        chunks = self.text_splitter.split_text(text)
        
        # 3. Store chunks in DB and prepare for FAISS
        faiss_texts = []
        faiss_metadatas = []
        
        for i, chunk_text in enumerate(chunks):
            # Save to DB
            db_chunk = models.DocumentChunk(
                document_id=document.id,
                chunk_index=i,
                text_content=chunk_text
            )
            db.add(db_chunk)
            
            # Prepare for FAISS
            faiss_texts.append(chunk_text)
            faiss_metadatas.append({
                "document_id": document.id,
                "filename": document.filename,
                "chunk_index": i,
                "text": chunk_text
            })
            
        db.commit()
        
        # 4. Add to FAISS
        faiss_store.add_documents(faiss_texts, faiss_metadatas)
        return True

document_processor = DocumentProcessor()
